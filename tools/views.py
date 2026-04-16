import os
import uuid
import json
import zipfile
import subprocess
import platform

from django.shortcuts import render, redirect
from django.http import JsonResponse, FileResponse, Http404, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
from django.core.files.storage import FileSystemStorage
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.forms import UserCreationForm
from django.contrib import messages

# PDF Processing Libraries
import fitz  # PyMuPDF
from PyPDF2 import PdfReader, PdfWriter, PdfMerger


# ==========================================
# TOOLS DATA CONFIGURATION
# ==========================================
TOOLS_DATA = {
    'merge': {
        'id': 'merge', 'title': 'Merge PDF files', 
        'desc': 'Combine PDFs in the order you want with the easiest PDF merger available.',
        'action_url': '/api/merge'
    },
    'split': {
        'id': 'split', 'title': 'Split PDF file', 
        'desc': 'Separate one page or a whole set for easy conversion into independent PDF files.',
        'action_url': '/api/split'
    },
    'compress': {
        'id': 'compress', 'title': 'Compress PDF file', 
        'desc': 'Reduce file size while optimizing for maximal PDF quality.',
        'action_url': '/api/compress'
    },
    'pdf-to-word': {
        'id': 'pdf-to-word', 'title': 'PDF to Word', 
        'desc': 'Easily convert your PDF files into easy to edit DOCX documents.',
        'action_url': '/api/pdf-to-word'
    },
    'pdf-to-jpg': {
        'id': 'pdf-to-jpg', 'title': 'PDF to JPG', 
        'desc': 'Convert each PDF page into a high-quality JPG image.',
        'action_url': '/api/pdf-to-jpg'
    },
    'unlock-pdf': {
        'id': 'unlock-pdf', 'title': 'Unlock PDF', 
        'desc': 'Remove password security from your PDF, so you can use it however you want.',
        'action_url': '/api/unlock-pdf'
    },
    'word-to-pdf': {
        'id': 'word-to-pdf', 'title': 'Word to PDF', 
        'desc': 'Make DOC and DOCX files easy to read by converting them to PDF.',
        'action_url': '/api/word-to-pdf'
    },
    'protect-pdf': {
        'id': 'protect-pdf', 'title': 'Protect PDF', 
        'desc': 'Encrypt your PDF with a password to prevent unauthorized access.',
        'action_url': '/api/protect-pdf'
    },
    'rotate-pdf': {
        'id': 'rotate-pdf', 'title': 'Rotate PDF', 
        'desc': 'Rotate your PDF pages however you need them. You can even rotate multiple PDFs at once!',
        'action_url': '/api/rotate-pdf'
    },
    'pdf-to-powerpoint': {
        'id': 'pdf-to-powerpoint', 'title': 'PDF to PowerPoint', 
        'desc': 'Turn your PDF files into easy to edit PPTX slideshows.',
        'action_url': '/api/pdf-to-powerpoint'
    },
    'pdf-to-excel': {
        'id': 'pdf-to-excel', 'title': 'PDF to Excel', 
        'desc': 'Pull data straight from PDFs into Excel spreadsheets in a few short seconds.',
        'action_url': '/api/pdf-to-excel'
    },
    'jpg-to-pdf': {
        'id': 'jpg-to-pdf', 'title': 'JPG to PDF', 
        'desc': 'Convert JPG images to PDF in seconds. Easily adjust orientation and margins.',
        'action_url': '/api/jpg-to-pdf'
    },
    'excel-to-pdf': {
        'id': 'excel-to-pdf', 'title': 'Excel to PDF', 
        'desc': 'Make EXCEL spreadsheets easy to read by converting them to PDF.',
        'action_url': '/api/excel-to-pdf'
    },
    'watermark': {
        'id': 'watermark', 'title': 'Watermark PDF', 
        'desc': 'Stamp text over your PDF in seconds. Choose the position and typography.',
        'action_url': '/api/watermark'
    },
    'powerpoint-to-pdf': {
        'id': 'powerpoint-to-pdf', 'title': 'PowerPoint to PDF', 
        'desc': 'Make PPT and PPTX slideshows easy to view by converting them to PDF.',
        'action_url': '/api/powerpoint-to-pdf'
    },
    'organize-pdf': {
        'id': 'organize-pdf', 'title': 'Organize PDF', 
        'desc': 'Sort, add and delete PDF pages. Drag and drop the page thumbnails and sort them however you like.',
        'action_url': '/api/organize-pdf'
    },
    'html-to-pdf': {
        'id': 'html-to-pdf', 'title': 'HTML to PDF', 
        'desc': 'Convert webpages in HTML to PDF. Copy and paste the URL of the page you want to convert.',
        'action_url': '/api/html-to-pdf'
    },
    'pdf-to-pdfa': {
        'id': 'pdf-to-pdfa', 'title': 'PDF to PDF/A', 
        'desc': 'Transform your PDF to PDF/A, the ISO-standardized version of PDF for long-term archiving.',
        'action_url': '/api/pdf-to-pdfa'
    },
    'sign-pdf': {
        'id': 'sign-pdf', 'title': 'Sign PDF', 
        'desc': 'Sign yourself or request electronic signatures from others.',
        'action_url': '/api/sign-pdf'
    },
}

# ==========================================
# PAGE VIEWS
# ==========================================
def home(request):
    return render(request, 'index.html', {'tools': TOOLS_DATA})

def show_tool(request, tool_id):
    if tool_id in TOOLS_DATA:
        return render(request, f'{tool_id}.html', {'tool': TOOLS_DATA[tool_id]})
    raise Http404("Tool not found")


# ==========================================
# MASTER DOWNLOAD FUNCTION (100% FIXED)
# ==========================================
def download_file(request, filename):
    # Har tool direct MEDIA_ROOT mein file save karta hai
    file_path = os.path.join(settings.MEDIA_ROOT, filename)
    
    if os.path.exists(file_path):
        # FileResponse as_attachment force karega browser ko PDF file download karne ke liye (HTML nahi)
        response = FileResponse(open(file_path, 'rb'), as_attachment=True, filename=filename)
        return response
        
    print(f"File Missing at: {file_path}")
    raise Http404("Bhai file nahi mili!")


# ==========================================
# API: MERGE PDF
# ==========================================
@csrf_exempt
def process_merge(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('pdf_files')
            page_order_json = request.POST.get('page_order')

            if not files:
                return JsonResponse({'status': 'error', 'message': 'No files uploaded.'})

            new_doc = fitz.open()

            if page_order_json:
                page_order = json.loads(page_order_json)
                doc_dict = {}
                for f in files:
                    f.seek(0)
                    doc_dict[f.name] = fitz.open(stream=f.read(), filetype="pdf")
                
                for item in page_order:
                    file_name = item['file']
                    page_num = int(item['original_page']) - 1
                    if file_name in doc_dict:
                        doc = doc_dict[file_name]
                        if 0 <= page_num < len(doc):
                            new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                
                for doc in doc_dict.values():
                    doc.close()
            else:
                for f in files:
                    f.seek(0)
                    doc = fitz.open(stream=f.read(), filetype="pdf")
                    new_doc.insert_pdf(doc)
                    doc.close()

            output_filename = f"merged_{uuid.uuid4().hex[:8]}.pdf"
            output_path = os.path.join(settings.MEDIA_ROOT, output_filename)
            new_doc.save(output_path)
            new_doc.close()

            return JsonResponse({'status': 'success', 'download_url': f"/download/{output_filename}"})

        except Exception as e:
            print(f"BACKEND CRASH REPORT: {e}") 
            return JsonResponse({'status': 'error', 'message': str(e)})
    
    return JsonResponse({'status': 'error', 'message': 'Invalid request method.'})


# ==========================================
# API: ORGANIZE PDF
# ==========================================
@csrf_exempt
def process_organize_pdf(request):
    if request.method == 'POST':
        try:
            files = request.FILES.getlist('pdf_files')
            page_data_json = request.POST.get('page_data')

            print("DEBUG: files=", files)
            print("DEBUG: page_data_json=", page_data_json)

            if not files:
                return JsonResponse({'status': 'error', 'message': 'No files uploaded.'})

            new_doc = fitz.open()

            if page_data_json:
                page_data = json.loads(page_data_json)
                print("DEBUG: page_data parsed=", page_data)
                
                if not page_data:
                    return JsonResponse({'status': 'error', 'message': 'No pages selected to save. Please keep at least one page.'})
                
                # Single file for Organize PDF
                f = files[0]
                f.seek(0)
                doc = fitz.open(stream=f.read(), filetype="pdf")
                
                print("DEBUG: source doc pages=", len(doc))
                
                for item in page_data:
                    page_num = int(item['page']) - 1 
                    rotation = int(item.get('rotation', 0))
                    
                    if 0 <= page_num < len(doc):
                        new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                        new_page_index = len(new_doc) - 1
                        new_page = new_doc[new_page_index]
                        if rotation:
                            new_page.set_rotation((new_page.rotation + rotation) % 360)
                
                doc.close()
            else:
                for f in files:
                    f.seek(0)
                    doc = fitz.open(stream=f.read(), filetype="pdf")
                    new_doc.insert_pdf(doc)
                    doc.close()

            print("DEBUG: new doc pages=", len(new_doc))

            out_name = f"organized_{uuid.uuid4().hex[:8]}.pdf"
            out_path = os.path.join(settings.MEDIA_ROOT, out_name)
            new_doc.save(out_path)
            new_doc.close()

            return JsonResponse({'status': 'success', 'download_url': f"/download/{out_name}"})

        except Exception as e:
            print(f"BACKEND CRASH REPORT: {e}") 
            return JsonResponse({'status': 'error', 'message': str(e)})
    
    return JsonResponse({'status': 'error', 'message': 'Invalid request method.'})



# ==========================================
# API: COMPRESS
# ==========================================
@csrf_exempt
def process_compress(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files:
            return JsonResponse({"error": "No file uploaded"}, status=400)
        
        compression_level = request.POST.get('compression_level', 'recommended')
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        
        try:
            f = files[0]
            path = fs.save(f.name, f)
            abs_path = os.path.join(settings.MEDIA_ROOT, path)
            out = f"comp_{uuid.uuid4().hex[:8]}.pdf"
            out_path = os.path.join(settings.MEDIA_ROOT, out)
            
            original_size = os.path.getsize(abs_path)
            
            gs_path = r"C:\Program Files\gs\gs10.07.0\bin\gswin64c.exe"
            # ==========================================
            # COMPRESSION LOGIC
            # ==========================================
            
            # 1. Map the frontend tier to the Ghostscript preset
            if compression_level == 'extreme':
                gs_setting = "/screen"  
            elif compression_level == 'recommended':
                gs_setting = "/ebook"   
            else: 
                gs_setting = "/printer" 

            # 2. Run Ghostscript with the selected preset
            gs_cmd = [
                gs_path, 
                "-sDEVICE=pdfwrite", 
                "-dCompatibilityLevel=1.4",
                f"-dPDFSETTINGS={gs_setting}",  # Dynamic setting based on tier
                "-dNOPAUSE", 
                "-dQUIET", 
                "-dBATCH",
                f"-sOutputFile={out_path}", 
                abs_path
            ]
            
            subprocess.run(gs_cmd, check=True)
            # ==========================================
            # ==========================================
            
            # Get the ACTUAL new file size after compression
            new_size = os.path.getsize(out_path)
            
            # If the file somehow got bigger (happens rarely with already-compressed PDFs), 
            # just serve the original file instead of faking the math.
            if new_size >= original_size:
                # Overwrite the larger output with the smaller original
                import shutil
                shutil.copyfile(abs_path, out_path)
                new_size = original_size
                savings = 0
            else:
                savings = max(0, round(((original_size - new_size) / original_size) * 100))
                
            # Clean up the original uploaded file
            os.remove(abs_path)
            
            return JsonResponse({
                "status": "success", 
                "download_url": f"/download/{out}", 
                "savings": savings, 
                "original_size": original_size, 
                "new_size": new_size
            })
            
        except Exception as e: 
            return JsonResponse({"error": f"Compression Error: {str(e)}"}, status=500)

# ==========================================
# API: PDF TO WORD
# ==========================================
@csrf_exempt
def process_pdf_to_word(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        from pdf2docx import Converter
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        try:
            f = files[0]
            path = fs.save(f.name, f)
            abs_path = os.path.join(settings.MEDIA_ROOT, path)
            out = f"word_{uuid.uuid4().hex}.docx"
            out_path = os.path.join(settings.MEDIA_ROOT, out)
            cv = Converter(abs_path)
            cv.convert(out_path)
            cv.close()
            os.remove(abs_path)
            return JsonResponse({"status": "success", "download_url": f"/download/{out}"})
        except Exception as e: return JsonResponse({"error": str(e)}, status=500)


# ==========================================
# API: PDF TO JPG
# ==========================================
@csrf_exempt
def process_pdf_to_jpg(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        try:
            zip_name = f"images_{unique_id}.zip"
            zip_path = os.path.join(settings.MEDIA_ROOT, zip_name)
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for f in files:
                    path = fs.save(f.name, f)
                    abs_path = os.path.join(settings.MEDIA_ROOT, path)
                    
                    doc = fitz.open(abs_path)
                    for i in range(len(doc)):
                        page = doc.load_page(i)
                        pix = page.get_pixmap(dpi=150)
                        img_name = f"page_{i+1}_{uuid.uuid4().hex[:4]}.jpg"
                        img_path = os.path.join(settings.MEDIA_ROOT, img_name)
                        pix.save(img_path)
                        zipf.write(img_path, arcname=img_name)
                        os.remove(img_path)
                    
                    doc.close()
                    os.remove(abs_path)
            return JsonResponse({"status": "success", "download_url": f"/download/{zip_name}"})
        except Exception as e: return JsonResponse({"error": str(e)}, status=500)


# ==========================================
# API: WORD TO PDF
# ==========================================
@csrf_exempt
def process_word_to_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        import mammoth
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                with open(filepath, "rb") as docx_file:
                    result = mammoth.convert_to_html(docx_file)
                    html_content = result.value 
                
                out_name = f"{os.path.splitext(file.name)[0]}_{unique_id[:4]}.pdf"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                
                doc = fitz.open()
                page = doc.new_page()
                rect = fitz.Rect(50, 50, 550, 800)
                page.insert_htmlbox(rect, html_content)
                
                doc.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"makemypdfs_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: UNLOCK PDF
# ==========================================
@csrf_exempt
def process_unlock_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        password = request.POST.get('pdf_password', '')
        
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        
        try:
            file = files[0]
            filename = fs.save(file.name, file)
            filepath = os.path.join(settings.MEDIA_ROOT, filename)
            
            doc = fitz.open(filepath)
            
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    doc.close()
                    os.remove(filepath)
                    return JsonResponse({"error": "Incorrect password. Please try again."}, status=401)
            
            out_name = f"unlocked_{uuid.uuid4().hex[:4]}_{file.name}"
            out_path = os.path.join(settings.MEDIA_ROOT, out_name)
            
            doc.save(out_path)
            doc.close()
            os.remove(filepath)

            return JsonResponse({"status": "success", "download_url": f"/download/{out_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: PROTECT PDF
# ==========================================
@csrf_exempt
def process_protect_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        password = request.POST.get('pdf_password', '')
        
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        if not password: return JsonResponse({"error": "Please provide a password"}, status=400)
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        
        try:
            file = files[0]
            filename = fs.save(file.name, file)
            filepath = os.path.join(settings.MEDIA_ROOT, filename)
            
            doc = fitz.open(filepath)
            out_name = f"protected_{unique_id[:4]}_{file.name}"
            out_path = os.path.join(settings.MEDIA_ROOT, out_name)
            
            doc.save(
                out_path, 
                encryption=fitz.PDF_ENCRYPT_AES_256, 
                user_pw=password, 
                owner_pw=password
            )
            doc.close()
            os.remove(filepath)

            return JsonResponse({"status": "success", "download_url": f"/download/{out_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: ROTATE PDF
# ==========================================
@csrf_exempt
def process_rotate_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        angle = int(request.POST.get('rotation_angle', 90))
        
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                doc = fitz.open(filepath)
                for page in doc:
                    page.set_rotation((page.rotation + angle) % 360)
                
                out_name = f"rotated_{unique_id[:4]}_{file.name}"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                doc.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"rotated_package_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)

            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: PDF TO POWERPOINT
# ==========================================
@csrf_exempt
def process_pdf_to_powerpoint(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        from pptx import Presentation
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                prs = Presentation()
                blank_slide_layout = prs.slide_layouts[6] 
                
                doc = fitz.open(filepath)
                for page in doc:
                    pix = page.get_pixmap(dpi=150)
                    img_path = os.path.join(settings.MEDIA_ROOT, f"temp_{unique_id}_{page.number}.png")
                    pix.save(img_path)
                    
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    os.remove(img_path)
                
                out_name = f"{os.path.splitext(file.name)[0]}_{unique_id[:4]}.pptx"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                prs.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"presentations_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: PDF TO EXCEL
# ==========================================
@csrf_exempt
def process_pdf_to_excel(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        import pdfplumber
        import pandas as pd
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                out_name = f"{os.path.splitext(file.name)[0]}_{unique_id[:4]}.xlsx"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                
                with pdfplumber.open(filepath) as pdf:
                    with pd.ExcelWriter(out_path, engine='openpyxl') as writer:
                        table_found = False
                        for i, page in enumerate(pdf.pages):
                            tables = page.extract_tables()
                            for j, table in enumerate(tables):
                                if table:
                                    df = pd.DataFrame(table)
                                    sheet_name = f"Page_{i+1}_Table_{j+1}"
                                    df.to_excel(writer, sheet_name=sheet_name[:31], index=False, header=False)
                                    table_found = True
                        
                        if not table_found:
                            df = pd.DataFrame([["No tabular data found in this PDF."]])
                            df.to_excel(writer, sheet_name="Result", index=False, header=False)
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"spreadsheets_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: JPG TO PDF
# ==========================================
@csrf_exempt
def process_jpg_to_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No images uploaded"}, status=400)
        
        orientation = request.POST.get('ui_orientation', 'portrait')
        margin_type = request.POST.get('ui_margin', 'no_margin')
        sorted_names = request.POST.get('sorted_names', '') 
        
        if sorted_names:
            order_list = sorted_names.split(',')
            file_dict = {f.name: f for f in files}
            sorted_files = []
            for name in order_list:
                if name in file_dict:
                    sorted_files.append(file_dict[name])
            files = sorted_files

        w, h = (595.0, 842.0) if orientation == 'portrait' else (842.0, 595.0)
        margin = 0
        if margin_type == 'small': margin = 20
        elif margin_type == 'big': margin = 50
        
        unique_id = uuid.uuid4().hex
        
        try:
            doc = fitz.open()
            for file in files:
                img_data = file.read()
                page = doc.new_page(width=w, height=h)
                rect = fitz.Rect(margin, margin, w - margin, h - margin)
                page.insert_image(rect, stream=img_data)
                
            out_name = f"images_converted_{unique_id[:4]}.pdf"
            out_path = os.path.join(settings.MEDIA_ROOT, out_name)
            doc.save(out_path)
            doc.close()
            
            return JsonResponse({"status": "success", "download_url": f"/download/{out_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: EXCEL TO PDF
# ==========================================
@csrf_exempt
def process_excel_to_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        import pandas as pd
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                xls = pd.read_excel(filepath, sheet_name=None)
                doc = fitz.open()
                
                for sheet_name, df in xls.items():
                    df = df.fillna('')
                    html_table = df.to_html(index=False)
                    html_content = f"""
                    <div style="font-family: sans-serif; font-size: 12px; color: #333;">
                        <h2 style="color: #1d6f42; padding-bottom: 10px; border-bottom: 2px solid #1d6f42;">Sheet: {sheet_name}</h2>
                        {html_table.replace('<table border="1" class="dataframe">', '<table style="border-collapse: collapse; width: 100%; border: 1px solid #ddd;">')
                                   .replace('<th>', '<th style="background-color: #1d6f42; color: white; padding: 8px; border: 1px solid #ddd; text-align: left;">')
                                   .replace('<td>', '<td style="padding: 8px; border: 1px solid #ddd;">')}
                    </div>
                    """
                    
                    page_height = max(842.0, float(150 + len(df) * 35)) 
                    page_width = max(595.0, float(100 + len(df.columns) * 120))
                    
                    page = doc.new_page(width=page_width, height=page_height)
                    rect = fitz.Rect(40, 40, page_width - 40, page_height - 40)
                    page.insert_htmlbox(rect, html_content)
                
                out_name = f"{os.path.splitext(file.name)[0]}_{unique_id[:4]}.pdf"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                doc.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"pdfs_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: WATERMARK
# ==========================================
@csrf_exempt
def process_watermark(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        wm_text = request.POST.get('watermark_text', 'MakeMyPDFs')
        position = request.POST.get('position', 'center')
        intensity = request.POST.get('intensity', 'medium')
        
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        if not wm_text.strip(): wm_text = "MakeMyPDFs"
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        if intensity == 'light': color = (0.85, 0.85, 0.85)
        elif intensity == 'dark': color = (0.2, 0.2, 0.2)
        else: color = (0.6, 0.6, 0.6) 
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                doc = fitz.open(filepath)
                for page in doc:
                    rect = page.rect
                    fontsize = 50
                    text_length = fitz.get_text_length(wm_text, fontname="helv", fontsize=fontsize)
                    
                    x = (rect.width - text_length) / 2
                    if position == 'top': y = 80 + fontsize 
                    elif position == 'bottom': y = rect.height - 80 
                    else: y = rect.height / 2 
                    
                    p = fitz.Point(x, y)
                    page.insert_text(p, wm_text, fontname="helv", fontsize=fontsize, color=color)
                
                out_name = f"watermarked_{unique_id[:4]}.pdf"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                doc.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"watermarks_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: POWERPOINT TO PDF
# ==========================================
@csrf_exempt
def process_powerpoint_to_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        if platform.system() == 'Darwin':
            libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        else:
            libreoffice_path = 'libreoffice'
            
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                subprocess.run([
                    libreoffice_path, '--headless', '--convert-to', 'pdf', 
                    filepath, '--outdir', settings.MEDIA_ROOT
                ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                
                lo_out_name = f"{os.path.splitext(file.name)[0]}.pdf"
                lo_out_path = os.path.join(settings.MEDIA_ROOT, lo_out_name)
                
                out_name = f"{os.path.splitext(file.name)[0]}_{unique_id[:4]}.pdf"
                final_out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                
                if os.path.exists(lo_out_path):
                    os.rename(lo_out_path, final_out_path)
                    processed_paths.append((final_out_path, out_name))
                
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"presentations_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": f"Engine error: {str(e)}"}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: HTML TO PDF
# ==========================================
@csrf_exempt
def process_html_to_pdf(request):
    if request.method == 'POST':
        url = request.POST.get('website_url')
        if not url: return JsonResponse({"error": "Please provide a valid URL"}, status=400)
        
        if not url.startswith('http://') and not url.startswith('https://'):
            url = 'https://' + url
            
        import pdfkit
        unique_id = uuid.uuid4().hex
        out_name = f"webpage_{unique_id[:6]}.pdf"
        out_path = os.path.join(settings.MEDIA_ROOT, out_name)
        
        try:
            options = {
                'page-size': 'A4', 'margin-top': '0mm', 'margin-right': '0mm',
                'margin-bottom': '0mm', 'margin-left': '0mm', 'encoding': "UTF-8",
            }
            pdfkit.from_url(url, out_path, options=options)
            return JsonResponse({"status": "success", "download_url": f"/download/{out_name}"})
        except Exception as e:
            return JsonResponse({"error": f"Failed to convert URL: {str(e)}"}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: PDF TO PDF/A
# ==========================================
@csrf_exempt
def process_pdf_to_pdfa(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                out_name = f"{os.path.splitext(file.name)[0]}_PDFA_{unique_id[:4]}.pdf"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                
                gs_cmd = [
                    'gs', '-dPDFA', '-dBATCH', '-dNOPAUSE', 
                    '-sColorConversionStrategy=UseDeviceIndependentColor', 
                    '-sDEVICE=pdfwrite', '-dPDFACompatibilityPolicy=2', 
                    f'-sOutputFile={out_path}', filepath
                ]
                
                subprocess.run(gs_cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"pdfa_archive_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": f"Engine error: {str(e)}"}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: SIGN PDF
# ==========================================
@csrf_exempt
def process_sign_pdf(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        sign_text = request.POST.get('sign_text', 'Signed Document')
        position = request.POST.get('position', 'bottom_right')
        page_choice = request.POST.get('page_choice', 'last')
        
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        if not sign_text.strip(): sign_text = "Signed Document"
        
        full_signature = f"Digitally Signed by: {sign_text}"
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        processed_paths = []
        
        try:
            for file in files:
                filename = fs.save(file.name, file)
                filepath = os.path.join(settings.MEDIA_ROOT, filename)
                
                doc = fitz.open(filepath)
                if page_choice == 'first': target_pages = [0]
                elif page_choice == 'last': target_pages = [len(doc) - 1]
                else: target_pages = range(len(doc)) 
                
                for page_num in target_pages:
                    page = doc[page_num]
                    rect = page.rect
                    fontsize = 16
                    text_length = fitz.get_text_length(full_signature, fontname="ti-ro", fontsize=fontsize)
                    margin_x, margin_y = 40, 40
                    
                    if position == 'bottom_right':
                        x, y = rect.width - text_length - margin_x, rect.height - margin_y
                    elif position == 'bottom_left':
                        x, y = margin_x, rect.height - margin_y
                    elif position == 'top_right':
                        x, y = rect.width - text_length - margin_x, margin_y + fontsize
                    else: 
                        x, y = margin_x, margin_y + fontsize
                    
                    p = fitz.Point(x, y)
                    page.insert_text(p, full_signature, fontname="ti-ro", fontsize=fontsize, color=(0.1, 0.2, 0.6))
                
                out_name = f"signed_{unique_id[:4]}.pdf"
                out_path = os.path.join(settings.MEDIA_ROOT, out_name)
                doc.save(out_path)
                doc.close()
                
                processed_paths.append((out_path, out_name))
                os.remove(filepath)

            if len(processed_paths) == 1:
                final_name = processed_paths[0][1]
            else:
                final_name = f"signed_docs_{unique_id}.zip"
                final_path = os.path.join(settings.MEDIA_ROOT, final_name)
                with zipfile.ZipFile(final_path, 'w') as zipf:
                    for p, fname in processed_paths:
                        zipf.write(p, arcname=fname)
                        os.remove(p)
                
            return JsonResponse({"status": "success", "download_url": f"/download/{final_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# API: SPLIT
# ==========================================
@csrf_exempt
def process_split(request):
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        if not files: return JsonResponse({"error": "No file uploaded"}, status=400)
        
        split_mode = request.POST.get('split_mode', 'range')
        fs = FileSystemStorage(location=settings.MEDIA_ROOT)
        unique_id = uuid.uuid4().hex
        
        try:
            file = files[0]
            filename = fs.save(file.name, file)
            filepath = os.path.join(settings.MEDIA_ROOT, filename)
            
            doc = fitz.open(filepath)
            out_name = f"split_{unique_id[:4]}.pdf"
            out_path = os.path.join(settings.MEDIA_ROOT, out_name)
            
            new_doc = fitz.open()
            
            if split_mode == 'range':
                from_page = int(request.POST.get('from_page', 1)) - 1
                to_page = int(request.POST.get('to_page', len(doc))) - 1
                from_page = max(0, from_page)
                to_page = min(len(doc)-1, to_page)
                
                if from_page <= to_page:
                    new_doc.insert_pdf(doc, from_page=from_page, to_page=to_page)
            else:
                pages_str = request.POST.get('pages_to_extract', '')
                if pages_str:
                    page_nums = []
                    for part in pages_str.split(','):
                        part = part.strip()
                        if '-' in part:
                            start, end = map(int, part.split('-'))
                            page_nums.extend(range(start-1, end))
                        else:
                            page_nums.append(int(part)-1)
                    
                    for p_num in page_nums:
                        if 0 <= p_num < len(doc):
                            new_doc.insert_pdf(doc, from_page=p_num, to_page=p_num)
            
            if len(new_doc) > 0:
                new_doc.save(out_path)
            new_doc.close()
            doc.close()
            os.remove(filepath)
            
            return JsonResponse({"status": "success", "download_url": f"/download/{out_name}"})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
            
    return JsonResponse({"error": "Invalid method"}, status=400)


# ==========================================
# AUTH & DASHBOARD VIEWS
# ==========================================
def signup_page(request):
    if request.method == 'POST':
        form = UserCreationForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('login') 
    else:
        form = UserCreationForm()
    return render(request, 'signup.html', {'form': form})

def login_page(request):
    if request.method == 'POST':
        u = request.POST.get('username')
        p = request.POST.get('password')
        user = authenticate(request, username=u, password=p)
        
        if user is not None:
            login(request, user)
            return redirect('/')  
        else:
            messages.error(request, 'Invalid Username or Password!')
            
    return render(request, 'login.html')

def logout_user(request):
    logout(request)
    return redirect('login')

def dashboard_view(request):
    return render(request, 'dashboard.html')